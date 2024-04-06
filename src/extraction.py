from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import re

def process_pptx_files(filepath_list):
    file_dict = {}
    for file_path in filepath_list:
        with open(file_path, "rb") as f:
            prs = Presentation(f)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text_runs.append(run.text)
                    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        for subshape in shape.shapes:
                            if subshape.has_text_frame:
                                for paragraph in subshape.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        text_runs.append(run.text)
                    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text:
                                    text_runs.append(cell.text)
            file_dict[file_path] = " ".join(text_runs)
    return file_dict

def extract_information(text):
    date_signature_regex = re.compile(r'Date de signature de la convention Natixis (\d{1,2} [^\s]+ \d{4})')
    montant_fasep_regex = re.compile(r'Montant du FASEP\s*([0-9 ]+€)')
    montant_acompte_regex = re.compile(r'Montant et date de paiement de l\'acompte\s*(\d{1,2} [^\s]+ \d{4}) pour (\d+ [^\s]+\s?euros?)')
    avis_service_economique_regex = re.compile(r'Avis sur le versement intermédiaire\s*(Favorable|Défavorable)')

    date_signature_match = date_signature_regex.search(text)
    date_signature = date_signature_match.group(1) if date_signature_match else "Non trouvée"

    montant_fasep_match = montant_fasep_regex.search(text)
    montant_fasep = montant_fasep_match.group(1) if montant_fasep_match else "Non trouvé"

    montant_acompte_match = montant_acompte_regex.search(text)
    montant_acompte = montant_acompte_match.group(2) if montant_acompte_match else "Non trouvé"
    date_acompte = montant_acompte_match.group(1) if montant_acompte_match else "Non trouvée"

    avis_service_economique_match = avis_service_economique_regex.search(text)
    avis_service_economique = avis_service_economique_match.group(1) if avis_service_economique_match else "Non trouvé"

    return {
        "Date de signature de la convention Natixis": date_signature,
        "Montant du FASEP": montant_fasep,
        "Montant de l'acompte": montant_acompte,
        "Date de paiement de l'acompte": date_acompte,
        "Avis du service économique pour le premier terme intermédiaire": avis_service_economique
    }
